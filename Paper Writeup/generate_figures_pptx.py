#!/usr/bin/env python3
"""
Generate conceptual figures for PQC on C-V2X Mode 4 paper.
Output: Paper Writeup/conceptual_figures.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

# ── Color palette ──
BLUE = RGBColor(0x1F, 0x77, 0xB4)       # ECDSA
ORANGE = RGBColor(0xFF, 0x7F, 0x0E)     # Falcon-512
GRAY = RGBColor(0x7F, 0x7F, 0x7F)       # Dilithium-2 (infeasible)
DARK = RGBColor(0x33, 0x33, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF0, 0xF0)
MED_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
GREEN = RGBColor(0x2C, 0xA0, 0x2C)
RED = RGBColor(0xD6, 0x27, 0x28)
PURPLE = RGBColor(0x94, 0x67, 0xBD)
TEAL = RGBColor(0x17, 0xBE, 0xCF)
LIGHT_BLUE = RGBColor(0xAE, 0xC7, 0xE8)
LIGHT_ORANGE = RGBColor(0xFF, 0xBB, 0x78)
LIGHT_GREEN = RGBColor(0x98, 0xDF, 0x8A)
VERY_LIGHT_BLUE = RGBColor(0xDE, 0xEB, 0xF7)
VERY_LIGHT_ORANGE = RGBColor(0xFD, 0xE8, 0xD0)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)


def add_title(slide, text, subtitle=""):
    """Add a styled title bar at the top of a slide."""
    # Title background bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
        Inches(13.333), Inches(0.9)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(28)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.LEFT
    tf.margin_left = Inches(0.5)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    if subtitle:
        shape2 = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.9),
            Inches(13.333), Inches(0.45)
        )
        shape2.fill.solid()
        shape2.fill.fore_color.rgb = RGBColor(0x55, 0x55, 0x55)
        shape2.line.fill.background()
        tf2 = shape2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(0xDD, 0xDD, 0xDD)
        p2.font.italic = True
        p2.alignment = PP_ALIGN.LEFT
        tf2.margin_left = Inches(0.5)
        tf2.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_box(slide, left, top, width, height, text, fill_color, font_size=12,
            font_color=DARK, bold=False, border_color=None, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    """Add a colored box with text."""
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.color.rgb = RGBColor(0x99, 0x99, 0x99)
        shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = PP_ALIGN.CENTER
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=DARK, width=Pt(2)):
    """Add a line with arrow."""
    conn = slide.shapes.add_connector(
        1, x1, y1, x2, y2  # 1 = straight connector
    )
    conn.line.color.rgb = color
    conn.line.width = width
    # End arrow
    conn.line.end_marker_style = 2  # triangle


def add_text(slide, left, top, width, height, text, font_size=11,
             font_color=DARK, bold=False, alignment=PP_ALIGN.LEFT):
    """Add a text box."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.alignment = alignment
    return txBox


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1: Title Slide
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
bg = slide.background
fill = bg.fill
fill.solid()
fill.fore_color.rgb = WHITE

# Large title
add_text(slide, Inches(1), Inches(2), Inches(11), Inches(1.5),
         "Conceptual Figures for Paper",
         font_size=36, bold=True, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(3.5), Inches(11), Inches(1),
         "Post-Quantum Cryptographic Authentication\non C-V2X Mode 4 Sidelink",
         font_size=22, font_color=RGBColor(0x55, 0x55, 0x55), alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(5), Inches(11), Inches(0.6),
         "Draft / Conceptual — For Professor Review",
         font_size=16, font_color=ORANGE, bold=True, alignment=PP_ALIGN.CENTER)

add_text(slide, Inches(1), Inches(6), Inches(11), Inches(0.5),
         "6 figures: System Architecture, SPDU Structure, SPS Flowchart, "
         "Intersection Layout, TBS Feasibility, Subchannel Grid",
         font_size=12, font_color=GRAY, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 2: System Architecture Diagram
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Fig. 1 — System Architecture",
          "Simulation stack: SUMO ↔ Veins ↔ openCV2X (OMNeT++) with PQC integration")

# Three main simulator boxes across the top
# SUMO
add_box(slide, Inches(0.5), Inches(1.8), Inches(3), Inches(1.2),
        "SUMO\n(Microscopic Traffic Simulator)", LIGHT_GREEN, 16, bold=True,
        border_color=GREEN)
add_text(slide, Inches(0.5), Inches(3.05), Inches(3), Inches(0.7),
         "• Vehicle mobility\n• Intersection signal control\n• 4-leg intersection scenario",
         font_size=10)

# Veins
add_box(slide, Inches(5.15), Inches(1.8), Inches(3), Inches(1.2),
        "Veins\n(V2X Simulation Framework)", LIGHT_BLUE, 16, bold=True,
        border_color=BLUE)
add_text(slide, Inches(5.15), Inches(3.05), Inches(3), Inches(0.7),
         "• TraCI interface to SUMO\n• Channel models (LOS/NLOS)\n• Mobility coupling",
         font_size=10)

# OMNeT++ / openCV2X
add_box(slide, Inches(9.8), Inches(1.8), Inches(3), Inches(1.2),
        "OMNeT++ / openCV2X\n(Network Simulator)", LIGHT_ORANGE, 16, bold=True,
        border_color=ORANGE)
add_text(slide, Inches(9.8), Inches(3.05), Inches(3), Inches(0.7),
         "• LTE sidelink (PC5)\n• Mode 4 SPS MAC\n• Full protocol stack",
         font_size=10)

# Arrows between simulators
add_arrow(slide, Inches(3.5), Inches(2.4), Inches(5.15), Inches(2.4), DARK)
add_arrow(slide, Inches(5.15), Inches(2.5), Inches(3.5), Inches(2.5), DARK)
add_text(slide, Inches(3.6), Inches(2.0), Inches(1.5), Inches(0.3),
         "TraCI", font_size=9, font_color=GRAY, alignment=PP_ALIGN.CENTER)

add_arrow(slide, Inches(8.15), Inches(2.4), Inches(9.8), Inches(2.4), DARK)
add_arrow(slide, Inches(9.8), Inches(2.5), Inches(8.15), Inches(2.5), DARK)
add_text(slide, Inches(8.3), Inches(2.0), Inches(1.3), Inches(0.3),
         "INET", font_size=9, font_color=GRAY, alignment=PP_ALIGN.CENTER)

# Protocol stack (below openCV2X, centered)
stack_left = Inches(3.5)
stack_w = Inches(6.3)
layer_h = Inches(0.65)

# Application layer
y = Inches(4.2)
add_box(slide, stack_left, y, stack_w, layer_h,
        "Application Layer (Mode4App)", VERY_LIGHT_ORANGE, 13, bold=True,
        border_color=ORANGE, shape_type=MSO_SHAPE.RECTANGLE)

# PQC module (inset in app layer)
add_box(slide, Inches(10.0), y + Inches(0.05), Inches(2.5), layer_h - Inches(0.1),
        "PQC Module\n(liboqs: Falcon-512,\nDilithium-2, ECDSA)", WHITE, 9, bold=False,
        border_color=RED, shape_type=MSO_SHAPE.RECTANGLE)

# IEEE 1609.2 Security
y += layer_h + Inches(0.08)
add_box(slide, stack_left, y, stack_w, layer_h * 0.8,
        "IEEE 1609.2 Security (SPDU: Header + BSM + Signature + Certificate/Digest)",
        RGBColor(0xFF, 0xEB, 0xCD), 11, bold=False,
        border_color=ORANGE, shape_type=MSO_SHAPE.RECTANGLE)

# PDCP
y += layer_h * 0.8 + Inches(0.08)
add_box(slide, stack_left, y, stack_w, layer_h * 0.65,
        "PDCP / RLC", VERY_LIGHT_BLUE, 12,
        border_color=BLUE, shape_type=MSO_SHAPE.RECTANGLE)

# MAC
y += layer_h * 0.65 + Inches(0.08)
add_box(slide, stack_left, y, stack_w, layer_h,
        "MAC Layer — Mode 4 SPS (LteMacVUeMode4)", LIGHT_BLUE, 13, bold=True,
        border_color=BLUE, shape_type=MSO_SHAPE.RECTANGLE)

# SPS details inside MAC
add_box(slide, Inches(10.0), y + Inches(0.05), Inches(2.5), layer_h - Inches(0.1),
        "SPS: Sensing + Selection\nReselection Counter\nCBR Congestion Ctrl", WHITE, 9,
        border_color=BLUE, shape_type=MSO_SHAPE.RECTANGLE)

# PHY
y += layer_h + Inches(0.08)
add_box(slide, stack_left, y, stack_w, layer_h * 0.8,
        "PHY — LTE Sidelink PC5 (PSCCH + PSSCH)", RGBColor(0xD5, 0xE8, 0xD4), 12,
        border_color=GREEN, shape_type=MSO_SHAPE.RECTANGLE)

# Channel
y += layer_h * 0.8 + Inches(0.08)
add_box(slide, stack_left, y, stack_w * 0.5, layer_h * 0.6,
        "LOS: Analytical\n(Dual-slope)", RGBColor(0xE8, 0xF5, 0xE9), 10,
        border_color=GREEN, shape_type=MSO_SHAPE.RECTANGLE)
add_box(slide, stack_left + stack_w * 0.52, y, stack_w * 0.48, layer_h * 0.6,
        "NLOS: Urban Microcell\n+ Nakagami Fading", RGBColor(0xFC, 0xE4, 0xEC), 10,
        border_color=RED, shape_type=MSO_SHAPE.RECTANGLE)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 3: IEEE 1609.2 SPDU Packet Structure
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Fig. 2 — IEEE 1609.2 SPDU Byte-Level Structure",
          "Comparison of SPDU sizes: ECDSA (compact) vs Falcon-512 (viable) vs Dilithium-2 (infeasible)")

# Table-like visualization
bar_top_start = Inches(2.0)
bar_height = Inches(0.9)
bar_gap = Inches(0.3)
label_w = Inches(2.0)
max_bar_w = Inches(10.0)  # pixels for max size (Dilithium)

# SPDU sizes (bytes)
algos = [
    ("ECDSA-256", 163, [
        ("Hdrs\n28B", 28, RGBColor(0xBD, 0xBD, 0xBD)),
        ("BSM\n39B", 39, LIGHT_BLUE),
        ("Sig\n64B", 64, LIGHT_ORANGE),
        ("Cert\n32B", 32, RGBColor(0xFF, 0x99, 0x99)),
    ]),
    ("Falcon-512\n(certificate)", 1694, [
        ("Hdrs 28B", 28, RGBColor(0xBD, 0xBD, 0xBD)),
        ("BSM 39B", 39, LIGHT_BLUE),
        ("Sig ~660B", 660, LIGHT_ORANGE),
        ("Cert (pub 897B + meta) ~967B", 967, RGBColor(0xFF, 0x99, 0x99)),
    ]),
    ("Falcon-512\n(digest)", 735, [
        ("Hdrs 28B", 28, RGBColor(0xBD, 0xBD, 0xBD)),
        ("BSM 39B", 39, LIGHT_BLUE),
        ("Sig ~660B", 660, LIGHT_ORANGE),
        ("Digest 8B", 8, RGBColor(0xFF, 0xCC, 0xCC)),
    ]),
    ("Dilithium-2\n(INFEASIBLE)", 2495, [
        ("Hdrs 28B", 28, RGBColor(0xBD, 0xBD, 0xBD)),
        ("BSM 39B", 39, LIGHT_BLUE),
        ("Sig 2420B", 2420, LIGHT_ORANGE),
        ("Digest 8B", 8, RGBColor(0xFF, 0xCC, 0xCC)),
    ]),
]

max_size = 2500  # scale reference
bar_left_start = Inches(2.5)

for i, (name, total, segments) in enumerate(algos):
    y = bar_top_start + i * (bar_height + bar_gap)

    # Label
    is_infeasible = "INFEASIBLE" in name
    label_color = RED if is_infeasible else DARK
    add_text(slide, Inches(0.3), y, label_w, bar_height,
             name, font_size=13, font_color=label_color, bold=True,
             alignment=PP_ALIGN.RIGHT)

    # Draw segments
    x_pos = bar_left_start
    for seg_label, seg_bytes, seg_color in segments:
        seg_w = max_bar_w * (seg_bytes / max_size)
        if seg_w < Inches(0.15):
            seg_w = Inches(0.15)
        shape = add_box(slide, x_pos, y + Inches(0.05), seg_w,
                       bar_height - Inches(0.1), seg_label,
                       seg_color, 8, border_color=DARK,
                       shape_type=MSO_SHAPE.RECTANGLE)
        x_pos += seg_w + Inches(0.02)

    # Total size label
    add_text(slide, x_pos + Inches(0.1), y, Inches(1.2), bar_height,
             f"= {total} B", font_size=13, bold=True,
             font_color=RED if is_infeasible else DARK)

# TBS limit line
tbs_x = bar_left_start + max_bar_w * (2481 / max_size)
line = slide.shapes.add_connector(1, tbs_x, bar_top_start - Inches(0.3),
                                   tbs_x, bar_top_start + 4 * (bar_height + bar_gap))
line.line.color.rgb = RED
line.line.width = Pt(2.5)
line.line.dash_style = 2  # dashed
add_text(slide, tbs_x - Inches(0.8), bar_top_start - Inches(0.6), Inches(1.6), Inches(0.3),
         "Max TBS = 2,481 B", font_size=11, font_color=RED, bold=True,
         alignment=PP_ALIGN.CENTER)

# Legend
add_text(slide, Inches(3), Inches(6.5), Inches(7), Inches(0.4),
         "■ Headers   ■ BSM Payload   ■ Signature   ■ Certificate/Digest   "
         "--- Max TBS (MCS 11, 10 SC)",
         font_size=10, font_color=GRAY)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 4: C-V2X Mode 4 SPS Flowchart
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Fig. 3 — C-V2X Mode 4 Semi-Persistent Scheduling (SPS) Flowchart",
          "Resource selection and reservation procedure per 3GPP TS 36.321 / TS 36.213")

# Flowchart boxes - vertical flow
fc_left = Inches(4.2)
fc_w = Inches(5.0)
fc_h = Inches(0.65)
fc_gap = Inches(0.18)
start_y = Inches(1.7)

steps = [
    ("BSM Generated at Application Layer\n(every 100 ms RRI)", LIGHT_GREEN, GREEN),
    ("Sense Channel: Measure RSSI/RSRP\nover Sensing Window (1000 ms)", LIGHT_BLUE, BLUE),
    ("Build Candidate Resource Set (SA)\nExclude resources with RSRP > threshold", LIGHT_BLUE, BLUE),
    ("Randomly Select from Top 20%\nLowest Average RSSI Candidates", LIGHT_BLUE, BLUE),
    ("Create SPS Grant\n(Fix MCS + Subchannels for reservation)", VERY_LIGHT_ORANGE, ORANGE),
    ("Transmit SPDU\n(PSCCH + PSSCH on reserved resources)", LIGHT_ORANGE, ORANGE),
    ("Decrement Reselection Counter\n(initialized to [5, 15])", VERY_LIGHT_BLUE, BLUE),
]

decision_y = None
for i, (text, fill, border) in enumerate(steps):
    y = start_y + i * (fc_h + fc_gap)
    add_box(slide, fc_left, y, fc_w, fc_h, text, fill, 11,
            border_color=border, shape_type=MSO_SHAPE.RECTANGLE)
    if i < len(steps) - 1:
        # Arrow down
        add_arrow(slide, fc_left + fc_w / 2, y + fc_h,
                 fc_left + fc_w / 2, y + fc_h + fc_gap, DARK)
    if i == len(steps) - 1:
        decision_y = y + fc_h + fc_gap

# Decision diamond (approximated with rotated rectangle text)
dy = decision_y
add_box(slide, fc_left + Inches(0.5), dy, fc_w - Inches(1), fc_h,
        "Counter = 0?", RGBColor(0xFF, 0xF9, 0xC4), 13, bold=True,
        border_color=RGBColor(0xFF, 0xA0, 0x00), shape_type=MSO_SHAPE.DIAMOND)

# Arrow down to decision
add_arrow(slide, fc_left + fc_w / 2, dy - fc_gap,
         fc_left + fc_w / 2, dy, DARK)

# "Yes" path - loops back to step 2 (new selection)
add_text(slide, fc_left - Inches(1.5), dy + Inches(0.1), Inches(1.2), Inches(0.4),
         "Yes →\nReselect", font_size=10, font_color=RED, bold=True, alignment=PP_ALIGN.RIGHT)
# Arrow from decision left
add_arrow(slide, fc_left + Inches(0.5), dy + fc_h / 2,
         fc_left - Inches(0.2), dy + fc_h / 2, RED)

# "No" path - keep with probability p
add_text(slide, fc_left + fc_w + Inches(0.2), dy + Inches(0.1), Inches(2.5), Inches(0.4),
         "No → Keep with\nprob. pKeep (0.8)", font_size=10, font_color=GREEN, bold=True)
# Arrow from decision right
add_arrow(slide, fc_left + fc_w - Inches(0.5), dy + fc_h / 2,
         fc_left + fc_w + Inches(0.15), dy + fc_h / 2, GREEN)

# Loop-back annotation
add_text(slide, fc_left - Inches(2.3), Inches(3.5), Inches(2), Inches(0.5),
         "← New sensing\n    & selection", font_size=10, font_color=RED)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 5: Intersection Scenario Layout
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Fig. 4 — Intersection Scenario Layout",
          "4-leg signalized intersection with varying vehicle densities (LOS A/C/F)")

# Draw intersection schematic
cx, cy = Inches(6.666), Inches(4.2)  # center
road_w = Inches(0.8)  # road width
road_len = Inches(2.8)

# Horizontal road
slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    cx - road_len - road_w / 2, cy - road_w / 2,
    road_len * 2 + road_w, road_w
).fill.solid()
slide.shapes[-1].fill.fore_color.rgb = RGBColor(0x90, 0x90, 0x90)
slide.shapes[-1].line.fill.background()

# Vertical road
slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    cx - road_w / 2, cy - road_len - road_w / 2,
    road_w, road_len * 2 + road_w
).fill.solid()
slide.shapes[-1].fill.fore_color.rgb = RGBColor(0x90, 0x90, 0x90)
slide.shapes[-1].line.fill.background()

# Center intersection (darker)
slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    cx - road_w / 2, cy - road_w / 2,
    road_w, road_w
).fill.solid()
slide.shapes[-1].fill.fore_color.rgb = RGBColor(0x70, 0x70, 0x70)
slide.shapes[-1].line.fill.background()

# Center dashed lines (lane dividers) - represented by thin lines
for offset in [road_w * 0.25, -road_w * 0.25]:
    # Horizontal lane divider (left approach)
    ln = slide.shapes.add_connector(
        1, cx - road_len - road_w / 2, cy + offset,
        cx - road_w / 2, cy + offset
    )
    ln.line.color.rgb = RGBColor(0xFF, 0xFF, 0x00)
    ln.line.width = Pt(1)
    ln.line.dash_style = 3
    # Horizontal lane divider (right approach)
    ln = slide.shapes.add_connector(
        1, cx + road_w / 2, cy + offset,
        cx + road_len + road_w / 2, cy + offset
    )
    ln.line.color.rgb = RGBColor(0xFF, 0xFF, 0x00)
    ln.line.width = Pt(1)
    ln.line.dash_style = 3
    # Vertical lane divider (top approach)
    ln = slide.shapes.add_connector(
        1, cx + offset, cy - road_len - road_w / 2,
        cx + offset, cy - road_w / 2
    )
    ln.line.color.rgb = RGBColor(0xFF, 0xFF, 0x00)
    ln.line.width = Pt(1)
    ln.line.dash_style = 3
    # Vertical lane divider (bottom approach)
    ln = slide.shapes.add_connector(
        1, cx + offset, cy + road_w / 2,
        cx + offset, cy + road_len + road_w / 2
    )
    ln.line.color.rgb = RGBColor(0xFF, 0xFF, 0x00)
    ln.line.width = Pt(1)
    ln.line.dash_style = 3

# Vehicle dots (small rectangles) on approaches
car_w, car_h = Inches(0.12), Inches(0.20)
car_colors = [BLUE, GREEN, ORANGE, PURPLE]
# West approach cars
for j in range(4):
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        cx - road_len + Inches(0.3 + j * 0.45), cy - road_w * 0.15,
        car_w, car_h
    ).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = BLUE
    slide.shapes[-1].line.fill.background()

# East approach cars
for j in range(3):
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        cx + road_w / 2 + Inches(0.4 + j * 0.5), cy + road_w * 0.05,
        car_w, car_h
    ).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = GREEN
    slide.shapes[-1].line.fill.background()

# North approach cars
for j in range(3):
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        cx + road_w * 0.05, cy - road_len + Inches(0.3 + j * 0.5),
        car_h, car_w  # rotated
    ).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = ORANGE
    slide.shapes[-1].line.fill.background()

# South approach cars
for j in range(3):
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        cx - road_w * 0.18, cy + road_w / 2 + Inches(0.3 + j * 0.5),
        car_h, car_w
    ).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = PURPLE
    slide.shapes[-1].line.fill.background()

# Approach labels
add_text(slide, cx - road_len - road_w, cy - road_w - Inches(0.4),
         Inches(1.5), Inches(0.3), "West (W)", font_size=12, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, cx + road_len - Inches(0.5), cy - road_w - Inches(0.4),
         Inches(1.5), Inches(0.3), "East (E)", font_size=12, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, cx + road_w + Inches(0.1), cy - road_len - Inches(0.1),
         Inches(1.5), Inches(0.3), "North (N)", font_size=12, bold=True, alignment=PP_ALIGN.CENTER)
add_text(slide, cx + road_w + Inches(0.1), cy + road_len - Inches(0.2),
         Inches(1.5), Inches(0.3), "South (S)", font_size=12, bold=True, alignment=PP_ALIGN.CENTER)

# Communication range circle (dashed)
comm_range = Inches(2.5)
circle = slide.shapes.add_shape(
    MSO_SHAPE.OVAL,
    cx - comm_range, cy - comm_range,
    comm_range * 2, comm_range * 2
)
circle.fill.background()
circle.line.color.rgb = RED
circle.line.width = Pt(1.5)
circle.line.dash_style = 3
add_text(slide, cx + comm_range - Inches(1.3), cy - comm_range - Inches(0.35),
         Inches(1.5), Inches(0.3), "Comm. range\n(~300 m)", font_size=9, font_color=RED)

# Density scenarios table on the right
table_left = Inches(10)
add_text(slide, table_left, Inches(1.7), Inches(3), Inches(0.4),
         "Traffic Density Scenarios", font_size=14, bold=True)

scenarios = [
    ("LOS A\n(Free-flow)", "10 veh/km", "724 veh/h", LIGHT_GREEN),
    ("LOS C\n(Stable flow)", "50 veh/km", "3,622 veh/h", RGBColor(0xFF, 0xFF, 0xCC)),
    ("LOS F\n(Breakdown)", "100 veh/km", "7,243 veh/h", RGBColor(0xFF, 0xCC, 0xCC)),
]

for i, (los, density, flow, color) in enumerate(scenarios):
    sy = Inches(2.3) + i * Inches(1.4)
    add_box(slide, table_left, sy, Inches(2.8), Inches(1.1), "", color, 10,
            border_color=DARK, shape_type=MSO_SHAPE.RECTANGLE)
    add_text(slide, table_left + Inches(0.1), sy + Inches(0.05),
             Inches(2.6), Inches(0.35), los, font_size=12, bold=True, alignment=PP_ALIGN.CENTER)
    add_text(slide, table_left + Inches(0.1), sy + Inches(0.4),
             Inches(2.6), Inches(0.3), f"Density: {density}", font_size=10, alignment=PP_ALIGN.CENTER)
    add_text(slide, table_left + Inches(0.1), sy + Inches(0.7),
             Inches(2.6), Inches(0.3), f"Flow: {flow}/approach", font_size=10, alignment=PP_ALIGN.CENTER)

# Simulation parameters
add_text(slide, Inches(0.5), Inches(6.3), Inches(5), Inches(1),
         "Parameters: 20 MHz BW, 10 subchannels, MCS 5–11\n"
         "RRI = 100 ms, Tx power = 23 dBm, 5.9 GHz\n"
         "Turn split: 80% through / 10% left / 10% right",
         font_size=10, font_color=GRAY)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 6: TBS vs SPDU Feasibility
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Fig. 5 — Transport Block Size (TBS) Feasibility Analysis",
          "Can the SPDU fit within the sidelink TBS? (3GPP TS 36.213 lookup tables)")

# Bar chart: TBS at different MCS for 10 subchannels
chart_left = Inches(1)
chart_top = Inches(2.0)
chart_w = Inches(8)
chart_h = Inches(4.5)
max_tbs = 3000  # scale max

# Y-axis
add_text(slide, Inches(0.1), Inches(3.8), Inches(0.8), Inches(0.4),
         "Bytes", font_size=11, font_color=GRAY, bold=True, alignment=PP_ALIGN.CENTER)

# MCS groups
mcs_data = [
    ("MCS 5\n(QPSK)", [("2 SC (20 RBs)", 252), ("4 SC (38 RBs)", 504),
                        ("6 SC (58 RBs)", 776), ("10 SC (98 RBs)", 1288)]),
    ("MCS 7\n(QPSK)", [("2 SC", 328), ("4 SC", 664),
                        ("6 SC", 1032), ("10 SC", 1736)]),
    ("MCS 11\n(16QAM)", [("2 SC", 488), ("4 SC", 984),
                          ("6 SC", 1544), ("10 SC", 2481)]),
]

group_w = Inches(2.2)
bar_w = Inches(0.4)
group_gap = Inches(0.6)

for gi, (mcs_label, bars) in enumerate(mcs_data):
    gx = chart_left + gi * (group_w + group_gap)

    # MCS label
    add_text(slide, gx, chart_top + chart_h + Inches(0.1), group_w, Inches(0.4),
             mcs_label, font_size=11, bold=True, alignment=PP_ALIGN.CENTER)

    for bi, (bar_label, tbs_val) in enumerate(bars):
        bx = gx + bi * (bar_w + Inches(0.08))
        bar_h_px = chart_h * (tbs_val / max_tbs)
        by = chart_top + chart_h - bar_h_px

        color = LIGHT_BLUE if bi < 3 else BLUE
        add_box(slide, bx, by, bar_w, bar_h_px, "",
                color, 8, border_color=DARK, shape_type=MSO_SHAPE.RECTANGLE)

        # Value on top
        add_text(slide, bx - Inches(0.1), by - Inches(0.3), bar_w + Inches(0.2), Inches(0.3),
                 str(tbs_val), font_size=8, bold=True, alignment=PP_ALIGN.CENTER)

# SPDU size reference lines
spdu_lines = [
    ("ECDSA SPDU = 163 B", 163, GREEN),
    ("Falcon-512 Digest = 735 B", 735, BLUE),
    ("Falcon-512 Cert = 1,694 B", 1694, ORANGE),
    ("Dilithium-2 = 2,495 B", 2495, RED),
    ("Max TBS = 2,481 B", 2481, RED),
]

for label, val, color in spdu_lines:
    line_y = chart_top + chart_h - chart_h * (val / max_tbs)
    if line_y < chart_top:
        line_y = chart_top
    ln = slide.shapes.add_connector(
        1, chart_left - Inches(0.2), line_y,
        chart_left + 3 * (group_w + group_gap) - group_gap, line_y
    )
    ln.line.color.rgb = color
    ln.line.width = Pt(1.5)
    ln.line.dash_style = 3

# Legend on the right
legend_left = Inches(9.3)
legend_top = Inches(2.2)
for i, (label, val, color) in enumerate(spdu_lines):
    ly = legend_top + i * Inches(0.45)
    # Color swatch
    sw = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 legend_left, ly + Inches(0.05),
                                 Inches(0.5), Inches(0.04))
    sw.fill.solid()
    sw.fill.fore_color.rgb = color
    sw.line.fill.background()
    add_text(slide, legend_left + Inches(0.6), ly, Inches(3), Inches(0.35),
             label, font_size=10, font_color=color, bold=("Max TBS" in label or "Dilithium" in label))

# Conclusion box
add_box(slide, Inches(9.3), Inches(5.0), Inches(3.5), Inches(1.5),
        "Dilithium-2: INFEASIBLE\n(2,495 B > 2,481 B max TBS)\n\n"
        "Falcon-512: Requires\nMCS ≥ 10, ≥ 8 subchannels\nfor cert-bearing SPDU",
        RGBColor(0xFF, 0xF0, 0xF0), 11, border_color=RED,
        shape_type=MSO_SHAPE.RECTANGLE)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 7: Subchannel Resource Grid
# ═══════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_title(slide, "Fig. 6 — Sidelink Subchannel Allocation per Algorithm",
          "Time-frequency resource grid showing PSCCH + PSSCH occupancy per transmission")

# Three grids side by side: ECDSA, Falcon-512 Digest, Falcon-512 Cert
grid_configs = [
    ("ECDSA-256\n(163 B)", 2, LIGHT_BLUE, BLUE),
    ("Falcon-512\nDigest (735 B)", 4, LIGHT_ORANGE, ORANGE),
    ("Falcon-512\nCert (1,694 B)", 8, RGBColor(0xFF, 0xCC, 0xCC), RED),
]

num_sc = 10
grid_left_start = Inches(0.8)
grid_w = Inches(3.5)
grid_gap = Inches(0.6)
grid_top = Inches(2.2)
cell_h = Inches(0.38)
cell_w = Inches(0.32)

for gi, (algo_label, sc_used, fill, border) in enumerate(grid_configs):
    gx = grid_left_start + gi * (grid_w + grid_gap)

    # Algorithm label
    add_text(slide, gx, Inches(1.6), grid_w, Inches(0.6),
             algo_label, font_size=14, bold=True, alignment=PP_ALIGN.CENTER,
             font_color=border)

    # Column headers (subchannel numbers)
    for sc in range(num_sc):
        add_text(slide, gx + Inches(0.5) + sc * cell_w, grid_top - Inches(0.25),
                 cell_w, Inches(0.25),
                 f"SC{sc + 1}", font_size=7, alignment=PP_ALIGN.CENTER, font_color=GRAY)

    # Row label
    add_text(slide, gx - Inches(0.1), grid_top, Inches(0.6), cell_h,
             "PSCCH", font_size=8, font_color=GRAY, alignment=PP_ALIGN.RIGHT)
    add_text(slide, gx - Inches(0.1), grid_top + cell_h + Inches(0.02), Inches(0.6), cell_h,
             "PSSCH", font_size=8, font_color=GRAY, alignment=PP_ALIGN.RIGHT)

    # Grid cells
    for row in range(2):  # PSCCH, PSSCH
        for sc in range(num_sc):
            cx_pos = gx + Inches(0.5) + sc * cell_w
            cy_pos = grid_top + row * (cell_h + Inches(0.02))

            if sc < sc_used:
                if row == 0:
                    # PSCCH (control) - first 2 RBs of first SC only
                    c = RGBColor(0xA5, 0xD6, 0xA7) if sc < 1 else fill
                    lbl = "SCI" if sc < 1 else ""
                else:
                    c = fill
                    lbl = "DATA" if sc == sc_used // 2 else ""
            else:
                c = LIGHT_GRAY
                lbl = ""

            shape = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, cx_pos, cy_pos, cell_w - Inches(0.02), cell_h
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = c
            shape.line.color.rgb = MED_GRAY
            shape.line.width = Pt(0.5)
            if lbl:
                tf = shape.text_frame
                tf.word_wrap = True
                tf.vertical_anchor = MSO_ANCHOR.MIDDLE
                p = tf.paragraphs[0]
                p.text = lbl
                p.font.size = Pt(7)
                p.font.color.rgb = DARK
                p.alignment = PP_ALIGN.CENTER

    # Subchannel count label
    add_text(slide, gx + Inches(0.3), grid_top + 2 * cell_h + Inches(0.15),
             Inches(2.5), Inches(0.3),
             f"{sc_used}/10 subchannels used ({sc_used * 10}%)",
             font_size=11, bold=True, font_color=border, alignment=PP_ALIGN.CENTER)

# Legend at bottom
add_text(slide, Inches(1), Inches(5.2), Inches(11), Inches(0.3),
         "■ SCI (Sidelink Control)   ■ PSSCH Data   □ Unused/Available   "
         "| Each subchannel = 10 RBs | 20 MHz bandwidth, adjacent scheme",
         font_size=10, font_color=GRAY, alignment=PP_ALIGN.CENTER)

# Impact note
add_box(slide, Inches(2), Inches(5.8), Inches(9), Inches(1.2),
        "Key Insight: Falcon-512 with full certificate consumes 80% of sidelink resources per TX,\n"
        "leaving only 2 subchannels for other vehicles in the same subframe.\n"
        "ECDSA uses only 20%, enabling 4× more concurrent transmissions.\n"
        "Under SPS, this allocation is FIXED for the entire grant duration — "
        "digest packets still reserve the same resources.",
        RGBColor(0xFF, 0xF8, 0xE1), 11, border_color=ORANGE,
        shape_type=MSO_SHAPE.RECTANGLE)


# ═══════════════════════════════════════════════════════════════════
# Save
# ═══════════════════════════════════════════════════════════════════
output_path = "/home/veins/src/simulte/Paper Writeup/conceptual_figures.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
