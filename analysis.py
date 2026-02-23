import os
import re
from pathlib import Path

import pandas as pd
import matplotlib.pyplot as plt

from pptx import Presentation
from pptx.util import Inches


# ========= USER SETTINGS =========
CSV_PATH = "results.csv"
SUMMARY_PATH = "sender_summary.csv"
OUTPUT_DIR = Path("plots")
IMAGE_EXT = ".png"
PPTX_PATH = OUTPUT_DIR / "plots.pptx"
# =================================


def sanitize_filename(name: str) -> str:
    return re.sub(r"[^A-Za-z0-9_.-]+", "_", str(name))


def sender_to_carname(sender: str) -> str:
    """Convert sender like carNonIp[3] → Car 4 (index + 1)."""
    m = re.search(r"\[(\d+)\]", str(sender))
    if m:
        num = int(m.group(1)) + 1
        return f"Car {num}"
    return str(sender)


def clear_output_folder(folder: Path, ext: str = ".png") -> None:
    """Clear PNGs but keep pptx."""
    folder.mkdir(parents=True, exist_ok=True)
    for f in folder.glob(f"*{ext}"):
        f.unlink()


def load_sender_summary(summary_path: str) -> dict:
    """Load sender_summary.csv into {sender: total_sent} dict."""
    if not Path(summary_path).exists():
        print(f"WARNING: {summary_path} not found, falling back to msgId-based sent estimate")
        return {}
    df_summary = pd.read_csv(summary_path, encoding="utf-8-sig", sep=r"[,\t]+", engine="python")
    df_summary.columns = df_summary.columns.str.strip()
    if "sender" not in df_summary.columns or "total_sent" not in df_summary.columns:
        print(f"WARNING: {summary_path} missing expected columns (found: {list(df_summary.columns)}), "
              f"falling back to msgId-based sent estimate")
        return {}
    df_summary["sender"] = df_summary["sender"].str.strip()
    return dict(zip(df_summary["sender"], df_summary["total_sent"]))


def compute_pdr(df, sent_counts: dict = None):
    """Compute overall Packet Delivery Ratio per sender.

    If sent_counts dict is provided, use ground-truth sent counts.
    Otherwise fall back to msgId max - min + 1 estimate.
    """
    pdr_results = {}

    for sender, group in df.groupby("sender"):
        received = len(group)

        if sent_counts and sender in sent_counts:
            sent = sent_counts[sender]
        else:
            msg_min = group["msgId"].min()
            msg_max = group["msgId"].max()
            sent = msg_max - msg_min + 1

        pdr = received / sent if sent > 0 else 0

        pdr_results[sender] = {
            "PDR": pdr,
            "Sent": sent,
            "Received": received
        }

    return pdr_results


def create_pptx_from_images(image_paths, pptx_path: Path) -> None:
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # blank slide

    for img_path in image_paths:
        slide = prs.slides.add_slide(blank_layout)
        left = Inches(0.5)
        top = Inches(0.5)
        width = prs.slide_width - Inches(1.0)
        slide.shapes.add_picture(str(img_path), left, top, width=width)

    prs.save(str(pptx_path))
    print(f"Saved PowerPoint to: {pptx_path}")


def main():
    clear_output_folder(OUTPUT_DIR, IMAGE_EXT)

    df = pd.read_csv(CSV_PATH)

    required_cols = ["t", "sender", "msgId", "dist_m",
                     "delay_ms", "Numer of Vehicles", "Algorithm", "spdu_size"]
    for col in required_cols:
        if col not in df.columns:
            raise ValueError(f"Missing column: {col}")

    # Numeric conversions
    df["t"] = pd.to_numeric(df["t"], errors="coerce")
    df["msgId"] = pd.to_numeric(df["msgId"], errors="coerce")
    df["dist_m"] = pd.to_numeric(df["dist_m"], errors="coerce")
    df["delay_ms"] = pd.to_numeric(df["delay_ms"], errors="coerce")
    df["Numer of Vehicles"] = pd.to_numeric(df["Numer of Vehicles"], errors="coerce")
    df["spdu_size"] = pd.to_numeric(df["spdu_size"], errors="coerce")

    df = df.dropna()

    if df.empty:
        print("No valid data after cleaning.")
        return

    # Extract global info
    algo_name = str(df.iloc[0]["Algorithm"])
    spdu_size = int(df.iloc[0]["spdu_size"])

    # Load ground-truth sent counts (if available)
    sent_counts = load_sender_summary(SUMMARY_PATH)

    # Compute overall PDR per sender (used in suptitle)
    pdr_dict = compute_pdr(df, sent_counts)

    # Weighted overall PDR: total_received / total_sent
    total_received = sum(v["Received"] for v in pdr_dict.values())
    total_sent = sum(v["Sent"] for v in pdr_dict.values())
    weighted_pdr = total_received / total_sent if total_sent > 0 else 0
    print(f"Weighted Overall PDR: {weighted_pdr * 100:.2f}% "
          f"({total_received} received / {total_sent} sent)")

    # Global axis limits
    global_x_min = 0
    global_x_max = df["t"].max() + 10

    global_y_min = 0
    dist_y_max = df["dist_m"].max() + 10
    delay_y_max = df["delay_ms"].max() + 10

    vehicles_y_min = 0
    vehicles_y_max = df["Numer of Vehicles"].max() + 1

    # Sort and deduplicate for the global vehicles-vs-time plot
    df_sorted = df.sort_values("t")
    df_vehicles = df_sorted.drop_duplicates(subset="t", keep="first")
    t_all = df_vehicles["t"].values
    vehicles_all = df_vehicles["Numer of Vehicles"].values

    # Keep track of images to add to PPTX
    image_paths = []

    # Loop by sender
    for sender, group in df.groupby("sender"):

        # Sort by time for per-sender analysis
        group = group.sort_values("t").copy()

        t = group["t"].values
        dist = group["dist_m"].values
        delay = group["delay_ms"].values
        car_name = sender_to_carname(sender)

        # ---------- PDR over time (cumulative) for this sender ----------
        # received_count: 1, 2, 3, ... as time progresses
        group["received_count"] = range(1, len(group) + 1)

        # msgId starts at 0, so max_msgId_so_far + 1 = exact packets sent up to that point
        group["max_msgId_so_far"] = group["msgId"].cummax()
        group["sent_est"] = group["max_msgId_so_far"] + 1
        group["sent_est"] = group["sent_est"].where(group["sent_est"] > 0, 1)

        group["pdr_time"] = group["received_count"] / group["sent_est"]
        # optional clip to [0,1]
        group["pdr_time"] = group["pdr_time"].clip(0, 1)

        t_pdr = group["t"].values
        pdr_time_vals = group["pdr_time"].values

        # Fetch overall PDR for title
        pdr_value = pdr_dict[sender]["PDR"]
        pdr_percent = f"{pdr_value * 100:.2f}%"

        # Create figure (now 4 subplots)
        fig, axes = plt.subplots(4, 1, figsize=(10, 11))

        # Full title including PDR
        fig.suptitle(
            f"Plot for {car_name}. Algorithm: {algo_name}, "
            f"SPDU Size: {spdu_size} Bytes, Overall PDR: {pdr_percent}",
            fontsize=14,
            fontweight="bold"
        )

        # -------------------------------------------------------
        # Plot 1 - Vehicles vs Time (global)
        # -------------------------------------------------------
        axes[0].plot(t_all, vehicles_all, marker="o", linestyle="-", label="Vehicles")
        axes[0].set_title("Number of Vehicles vs Time")
        axes[0].set_ylabel("Vehicles")
        axes[0].grid(True)
        axes[0].legend()
        axes[0].set_xlim(global_x_min, global_x_max)
        axes[0].set_ylim(vehicles_y_min, vehicles_y_max)

        # -------------------------------------------------------
        # Plot 2 - Distance vs Time (per car)
        # -------------------------------------------------------
        axes[1].plot(t, dist, marker="o", linestyle="-", label="Distance")
        axes[1].set_title("Distance vs Time")
        axes[1].set_ylabel("Distance (m)")
        axes[1].grid(True)
        axes[1].legend()
        axes[1].set_xlim(global_x_min, global_x_max)
        axes[1].set_ylim(global_y_min, dist_y_max)

        # -------------------------------------------------------
        # Plot 3 - Delay vs Time (per car)
        # -------------------------------------------------------
        axes[2].plot(t, delay, marker="o", linestyle="-", label="Delay")
        axes[2].set_title("Delay vs Time")
        axes[2].set_xlabel("Time (t)")
        axes[2].set_ylabel("Delay (ms)")
        axes[2].grid(True)
        axes[2].legend()
        axes[2].set_xlim(global_x_min, global_x_max)
        axes[2].set_ylim(global_y_min, delay_y_max)

        # -------------------------------------------------------
        # Plot 4 - Cumulative PDR vs Time (per car) - SCATTER ONLY
        # -------------------------------------------------------
        axes[3].scatter(t_pdr, pdr_time_vals, label="PDR (cumulative)", marker="o")
        axes[3].axhline(0.9, color="red", linestyle="--", linewidth=1.5, label="PDR = 0.9")
        axes[3].set_title("Cumulative PDR vs Time")
        axes[3].set_xlabel("Time (t)")
        axes[3].set_ylabel("PDR")
        axes[3].grid(True)
        axes[3].legend()
        axes[3].set_xlim(global_x_min, global_x_max)
        axes[3].set_ylim(0, 1.05)  # PDR between 0 and 1

        plt.tight_layout(rect=[0, 0, 1, 0.95])

        # Save figure
        fname = OUTPUT_DIR / f"{sanitize_filename(car_name)}{IMAGE_EXT}"
        plt.savefig(fname, dpi=300)
        plt.close(fig)

        image_paths.append(fname)

    # -------------------------------------------------------
    # Overall PDR vs Average Distance plot
    # -------------------------------------------------------
    sender_stats = []
    for sender, group in df.groupby("sender"):
        avg_dist = group["dist_m"].mean()
        info = pdr_dict[sender]
        sender_stats.append({
            "sender": sender,
            "car_name": sender_to_carname(sender),
            "avg_dist": avg_dist,
            "pdr": info["PDR"],
            "sent": info["Sent"],
            "received": info["Received"],
        })
    df_stats = pd.DataFrame(sender_stats).sort_values("avg_dist")

    fig, ax = plt.subplots(figsize=(10, 6))
    ax.scatter(df_stats["avg_dist"], df_stats["pdr"] * 100, marker="o", s=50)

    # Label each point with car name
    for _, row in df_stats.iterrows():
        ax.annotate(row["car_name"], (row["avg_dist"], row["pdr"] * 100),
                    textcoords="offset points", xytext=(5, 5), fontsize=7)

    ax.axhline(90, color="red", linestyle="--", linewidth=1.5, label="PDR = 90%")
    ax.set_title(
        f"Per-Vehicle PDR vs Average Distance from RSU\n"
        f"Algorithm: {algo_name}, SPDU: {spdu_size} B, "
        f"Weighted PDR: {weighted_pdr * 100:.2f}%",
        fontsize=13, fontweight="bold"
    )
    ax.set_xlabel("Average Distance from RSU (m)")
    ax.set_ylabel("PDR (%)")
    ax.set_ylim(0, 105)
    ax.grid(True)
    ax.legend()
    plt.tight_layout()

    fname = OUTPUT_DIR / f"PDR_vs_Distance{IMAGE_EXT}"
    plt.savefig(fname, dpi=300)
    plt.close(fig)
    image_paths.append(fname)
    print(f"Saved PDR vs Distance plot: {fname}")

    # -------------------------------------------------------
    # Binned PDR vs Distance plot
    # -------------------------------------------------------
    bin_width = 50  # meters
    max_dist = df_stats["avg_dist"].max()
    bins = list(range(0, int(max_dist) + bin_width + 1, bin_width))
    df_stats["dist_bin"] = pd.cut(df_stats["avg_dist"], bins=bins, right=False)

    binned = df_stats.groupby("dist_bin", observed=True).agg(
        avg_pdr=("pdr", "mean"),
        num_vehicles=("sender", "count"),
        total_sent=("sent", "sum"),
        total_received=("received", "sum"),
    ).reset_index()
    binned["weighted_pdr"] = binned["total_received"] / binned["total_sent"]
    binned["bin_center"] = binned["dist_bin"].apply(lambda x: x.mid)

    fig, ax = plt.subplots(figsize=(10, 6))
    bars = ax.bar(binned["bin_center"], binned["weighted_pdr"] * 100,
                  width=bin_width * 0.8, edgecolor="black", alpha=0.7)

    # Add vehicle count labels on bars
    for bar, n in zip(bars, binned["num_vehicles"]):
        ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height() + 1,
                f"n={n}", ha="center", va="bottom", fontsize=9)

    ax.axhline(90, color="red", linestyle="--", linewidth=1.5, label="PDR = 90%")
    ax.set_title(
        f"Weighted PDR vs Distance (binned, {bin_width}m intervals)\n"
        f"Algorithm: {algo_name}, SPDU: {spdu_size} B, "
        f"Weighted PDR: {weighted_pdr * 100:.2f}%",
        fontsize=13, fontweight="bold"
    )
    ax.set_xlabel("Distance from RSU (m)")
    ax.set_ylabel("PDR (%)")
    ax.set_ylim(0, 105)
    ax.grid(True, axis="y")
    ax.legend()
    plt.tight_layout()

    fname = OUTPUT_DIR / f"PDR_vs_Distance_Binned{IMAGE_EXT}"
    plt.savefig(fname, dpi=300)
    plt.close(fig)
    image_paths.append(fname)
    print(f"Saved binned PDR vs Distance plot: {fname}")

    # Create PPTX
    if image_paths:
        create_pptx_from_images(image_paths, PPTX_PATH)


if __name__ == "__main__":
    main()

